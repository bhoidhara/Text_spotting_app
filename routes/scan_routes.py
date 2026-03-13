import io
import os
from uuid import uuid4

from flask import (
    abort,
    current_app,
    flash,
    jsonify,
    redirect,
    render_template,
    request,
    send_file,
    send_from_directory,
    session,
    url_for,
)

from services.ocr import (
    auto_correct_text,
    clean_text,
    detect_intent,
    detect_language,
    extract_actions,
    ocr_image,
    student_pack,
)
from services.scans import delete_scan as delete_scan_record, get_scan, list_scans, log_export, log_translation, upsert_scan
from services.storage import delete_from_storage, download_from_storage, upload_to_storage
from services.supabase_client import require_storage_client, supabase_storage_client
from services.scans import _supabase_only
from utils.auth import get_user_id, require_login
from utils.helpers import allowed_file, now_iso, safe_int, safe_slug
from werkzeug.utils import secure_filename

try:
    from PIL import Image
except Exception:
    Image = None


def register_scan_routes(app):
    def _api_user_id():
        return request.headers.get("X-User-Id") or session.get("user_id")

    def _build_crop_box():
        crop_x = request.form.get("crop_x")
        crop_y = request.form.get("crop_y")
        crop_w = request.form.get("crop_w")
        crop_h = request.form.get("crop_h")
        if crop_x and crop_y and crop_w and crop_h:
            x = safe_int(crop_x)
            y = safe_int(crop_y)
            w = safe_int(crop_w)
            h = safe_int(crop_h)
            if w <= 1 or h <= 1:
                return None
            if x < 0:
                x = 0
            if y < 0:
                y = 0
            return (x, y, x + w, y + h)
        return None

    def _process_upload(files, user_id, use_flash=True):
        if not files or all(file.filename == "" for file in files):
            return None, [], "Please select at least one image."

        if Image is None:
            return None, [], "Pillow is not installed. OCR cannot run."

        lang = request.form.get("lang", "eng") or "eng"
        cleanup = request.form.get("cleanup") == "on"
        autocorrect = request.form.get("autocorrect") == "on"
        detect_intent_flag = request.form.get("detect_intent") == "on"
        student_mode = request.form.get("student_mode") == "on"
        advanced_ocr = request.form.get("advanced_ocr") == "on"
        fast_param = request.form.get("fast_ocr")
        if advanced_ocr:
            fast_ocr = False
        else:
            # Default to fast OCR for all uploads unless explicitly disabled.
            fast_ocr = fast_param != "off"
        privacy_mode = request.form.get("privacy_mode") == "on"
        skip_storage = request.form.get("skip_storage") == "on"
        crop_box = _build_crop_box()

        combined_text = []
        combined_conf = []
        combined_low_conf = []
        combined_line_conf = []
        image_paths = []
        page_number = 0
        scan_id = str(uuid4())
        safe_user = safe_slug(user_id)
        store_local = os.getenv("STORE_LOCAL_UPLOADS", "false").lower() not in {"false", "0", "no"}
        upload_root = current_app.config["UPLOAD_FOLDER"]
        warnings = []
        max_upload_mb = int(current_app.config.get("MAX_UPLOAD_MB", 40))
        soft_upload_mb = int(current_app.config.get("SOFT_UPLOAD_MB", 0))
        max_bytes = max_upload_mb * 1024 * 1024
        soft_bytes = soft_upload_mb * 1024 * 1024
        max_pdf_pages = int(current_app.config.get("MAX_PDF_PAGES", 0))
        max_image_side = int(current_app.config.get("MAX_IMAGE_SIDE", 3200))

        storage_env_on = os.getenv("SUPABASE_USE_STORAGE", "true").lower() not in {"false", "0", "no"}
        storage_available = False
        if storage_env_on and not privacy_mode and not skip_storage:
            try:
                require_storage_client()
                storage_available = supabase_storage_client() is not None
            except Exception:
                storage_available = False
                if _supabase_only():
                    return None, warnings, "Storage is not configured."
                warnings.append("Cloud storage not available; keeping files locally.")
        if skip_storage:
            warnings.append("Storage upload skipped for faster scanning.")


        image_exts = {".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".webp", ".heic", ".heif"}

        empty_pages = 0

        def _add_page(text, avg_conf=0.0, low_conf=None, line_conf=None):
            nonlocal page_number
            nonlocal empty_pages
            page_number += 1
            page_text = (text or "").strip()
            if not page_text:
                empty_pages += 1
                return
            page_block = f"--- Page {page_number} ---\n{page_text}"
            combined_text.append(page_block.strip())
            combined_conf.append(avg_conf)
            if low_conf:
                combined_low_conf.extend(low_conf)
            if line_conf:
                for item in line_conf:
                    combined_line_conf.append(
                        {
                            "page": page_number,
                            "text": item.get("text", ""),
                            "conf": item.get("conf", 0.0),
                        }
                    )

        for file in files:
            size_bytes = None
            try:
                if hasattr(file, "stream") and hasattr(file.stream, "seek"):
                    pos = file.stream.tell()
                    file.stream.seek(0, os.SEEK_END)
                    size_bytes = file.stream.tell()
                    file.stream.seek(pos or 0)
            except Exception:
                size_bytes = None
            if soft_upload_mb > 0 and size_bytes:
                if size_bytes > max_bytes:
                    warnings.append(
                        f"{file.filename} is very large. Auto-optimizing for speed."
                    )
                    advanced_ocr = False
                    fast_ocr = True
                elif size_bytes > soft_bytes:
                    warnings.append(
                        f"{file.filename} is large. Auto-optimizing for speed."
                    )
                    advanced_ocr = False
                    fast_ocr = True
            mime = (getattr(file, "mimetype", "") or "").lower()
            is_image_mime = mime.startswith("image/")
            is_text_mime = mime.startswith("text/") or mime in {
                "application/json",
                "application/xml",
                "application/rtf",
                "application/x-rtf",
                "text/rtf",
                "text/csv",
            }
            if not file or (
                not allowed_file(file.filename, current_app.config["ALLOWED_EXT"])
                and not is_image_mime
                and not is_text_mime
            ):
                warnings.append(f"Skipped unsupported file: {file.filename}")
                continue

            filename = f"{uuid4().hex}_{secure_filename(file.filename)}"
            local_dir = os.path.join(upload_root, safe_user, scan_id)
            os.makedirs(local_dir, exist_ok=True)
            local_rel_path = os.path.join(safe_user, scan_id, filename)
            file_path = os.path.join(upload_root, local_rel_path)
            file.save(file_path)

            _, ext = os.path.splitext(file.filename.lower())

            if ext in {".txt"} or is_text_mime:
                try:
                    with open(file_path, "r", encoding="utf-8", errors="ignore") as handle:
                        content = handle.read()
                except Exception as exc:
                    warnings.append(f"Failed to read {file.filename}: {exc}")
                    content = ""
                if content.strip():
                    warnings.append(f"Imported text from {file.filename} (no OCR).")
                    line_conf = [
                        {"text": line.strip(), "conf": 100.0}
                        for line in content.splitlines()
                        if line.strip()
                    ]
                    _add_page(content, avg_conf=100.0, low_conf=[], line_conf=line_conf)
                else:
                    warnings.append(f"Empty text file: {file.filename}")

                if privacy_mode and _supabase_only():
                    try:
                        os.remove(file_path)
                    except Exception:
                        pass
                    return None, warnings, "Privacy mode requires local storage, but local is disabled."
                if privacy_mode or not store_local:
                    try:
                        os.remove(file_path)
                    except Exception:
                        pass
                continue

            if ext in {".docx"}:
                try:
                    from docx import Document
                except Exception as exc:
                    warnings.append(f"DOCX support not installed: {exc}")
                    continue
                try:
                    doc = Document(file_path)
                    content = "\n".join(p.text for p in doc.paragraphs if p.text)
                except Exception as exc:
                    warnings.append(f"Failed to read {file.filename}: {exc}")
                    content = ""
                if content.strip():
                    warnings.append(f"Imported text from {file.filename} (no OCR).")
                    line_conf = [
                        {"text": line.strip(), "conf": 100.0}
                        for line in content.splitlines()
                        if line.strip()
                    ]
                    _add_page(content, avg_conf=100.0, low_conf=[], line_conf=line_conf)
                else:
                    warnings.append(f"Empty DOCX file: {file.filename}")

                if privacy_mode and _supabase_only():
                    try:
                        os.remove(file_path)
                    except Exception:
                        pass
                    return None, warnings, "Privacy mode requires local storage, but local is disabled."
                if privacy_mode or not store_local:
                    try:
                        os.remove(file_path)
                    except Exception:
                        pass
                continue

            if ext in {".pptx"}:
                try:
                    from pptx import Presentation
                except Exception as exc:
                    warnings.append(f"PPTX support not installed: {exc}")
                    continue
                try:
                    deck = Presentation(file_path)
                    slides_text = []
                    for slide in deck.slides:
                        parts = []
                        for shape in slide.shapes:
                            text = getattr(shape, "text", "") or ""
                            if text.strip():
                                parts.append(text.strip())
                        slides_text.append("\n".join(parts).strip())
                except Exception as exc:
                    warnings.append(f"Failed to read {file.filename}: {exc}")
                    slides_text = []
                if any(text.strip() for text in slides_text):
                    warnings.append(f"Imported text from {file.filename} (no OCR).")
                    for slide_text in slides_text:
                        if slide_text.strip():
                            line_conf = [
                                {"text": line.strip(), "conf": 100.0}
                                for line in slide_text.splitlines()
                                if line.strip()
                            ]
                            _add_page(slide_text, avg_conf=100.0, low_conf=[], line_conf=line_conf)
                else:
                    warnings.append(f"Empty PPTX file: {file.filename}")

                if privacy_mode and _supabase_only():
                    try:
                        os.remove(file_path)
                    except Exception:
                        pass
                    return None, warnings, "Privacy mode requires local storage, but local is disabled."
                if privacy_mode or not store_local:
                    try:
                        os.remove(file_path)
                    except Exception:
                        pass
                continue

            if ext in {".pdf"}:
                try:
                    from pdf2image import convert_from_path, pdfinfo_from_path
                except Exception as exc:
                    warnings.append(f"PDF support not installed: {exc}")
                    continue
                try:
                    base_dpi = 180 if advanced_ocr else 100 if fast_ocr else 150
                    first_page = 1
                    last_page = None
                    total_pages = 0
                    try:
                        info = pdfinfo_from_path(file_path)
                        total_pages = safe_int(info.get("Pages", 0), 0)
                    except Exception:
                        total_pages = 0

                    pdf_dpi = base_dpi
                    if total_pages >= 25:
                        pdf_dpi = max(80, base_dpi - 60)
                        warnings.append(
                            f"Large PDF detected ({total_pages} pages). Using lower DPI for speed."
                        )
                    elif total_pages >= 12:
                        pdf_dpi = max(90, base_dpi - 40)
                        warnings.append(
                            f"PDF has {total_pages} pages. Using optimized DPI for speed."
                        )

                    if max_pdf_pages > 0:
                        last_page = max_pdf_pages
                        if total_pages and total_pages <= max_pdf_pages:
                            last_page = total_pages
                        if total_pages and total_pages > max_pdf_pages:
                            warnings.append(
                                f"PDF has {total_pages} pages. Processing first {max_pdf_pages} pages for speed."
                            )
                        elif total_pages == 0:
                            warnings.append(
                                f"Processing first {max_pdf_pages} pages for speed."
                            )
                    pages = convert_from_path(
                        file_path, dpi=pdf_dpi, first_page=first_page, last_page=last_page
                    )
                except Exception as exc:
                    warnings.append(f"Failed to read PDF {file.filename}: {exc}")
                    pages = []

                for page in pages:
                    if crop_box:
                        try:
                            page = page.crop(crop_box)
                        except Exception:
                            pass
                    try:
                        page_text, avg_conf, low_conf, line_conf = ocr_image(
                            page, lang=lang, advanced=advanced_ocr, fast=fast_ocr and not advanced_ocr
                        )
                    except Exception as exc:
                        if not advanced_ocr:
                            try:
                                page_text, avg_conf, low_conf, line_conf = ocr_image(
                                    page, lang=lang, advanced=False, fast=False
                                )
                                warnings.append(f"Fallback OCR used for {file.filename}.")
                            except Exception as exc2:
                                if not fast_ocr:
                                    try:
                                        page_text, avg_conf, low_conf, line_conf = ocr_image(
                                            page, lang=lang, advanced=True, fast=False
                                        )
                                        warnings.append(f"Fallback OCR used for {file.filename}.")
                                    except Exception as exc3:
                                        warnings.append(f"OCR failed on {file.filename}: {exc3}")
                                        continue
                                else:
                                    warnings.append(f"OCR failed on {file.filename}: {exc2}")
                                    continue
                        else:
                            warnings.append(f"OCR failed on {file.filename}: {exc}")
                            continue
                    if fast_ocr and not advanced_ocr:
                        short_text = len((page_text or "").strip()) < 12
                        low_conf = avg_conf and avg_conf < 40
                        if short_text or low_conf:
                            try:
                                page_text, avg_conf, low_conf, line_conf = ocr_image(
                                    page, lang=lang, advanced=False, fast=False
                                )
                                if (page_text or "").strip():
                                    warnings.append(f"Fallback OCR used for {file.filename}.")
                            except Exception:
                                pass
                    _add_page(page_text, avg_conf=avg_conf, low_conf=low_conf, line_conf=line_conf)
                    if not (page_text or "").strip() and fast_ocr and not advanced_ocr:
                        try:
                            page_text, avg_conf, low_conf, line_conf = ocr_image(
                                page, lang=lang, advanced=False, fast=False
                            )
                            if (page_text or "").strip():
                                warnings.append(f"Fallback OCR used for {file.filename}.")
                                _add_page(page_text, avg_conf=avg_conf, low_conf=low_conf, line_conf=line_conf)
                        except Exception:
                            pass
                    if not (page_text or "").strip() and not advanced_ocr and not fast_ocr:
                        try:
                            page_text, avg_conf, low_conf, line_conf = ocr_image(
                                page, lang=lang, advanced=True, fast=False
                            )
                            if (page_text or "").strip():
                                warnings.append(f"Fallback OCR used for {file.filename}.")
                                _add_page(page_text, avg_conf=avg_conf, low_conf=low_conf, line_conf=line_conf)
                        except Exception:
                            pass

                if privacy_mode and _supabase_only():
                    try:
                        os.remove(file_path)
                    except Exception:
                        pass
                    return None, warnings, "Privacy mode requires local storage, but local is disabled."
                if privacy_mode or not store_local:
                    try:
                        os.remove(file_path)
                    except Exception:
                        pass
                continue

            try:
                image = Image.open(file_path)
            except Exception:
                warnings.append(f"Failed to open {file.filename}")
                continue

            if max_image_side and (
                image.width > max_image_side or image.height > max_image_side
            ):
                try:
                    image.thumbnail(
                        (max_image_side, max_image_side),
                        resample=getattr(image, "LANCZOS", 1),
                    )
                    warnings.append(
                        f"Large image scaled down for faster OCR: {file.filename}"
                    )
                except Exception:
                    pass

            if crop_box:
                try:
                    image = image.crop(crop_box)
                except Exception:
                    pass

            try:
                page_text, avg_conf, low_conf, line_conf = ocr_image(
                    image, lang=lang, advanced=advanced_ocr, fast=fast_ocr and not advanced_ocr
                )
            except Exception as exc:
                if not advanced_ocr:
                    try:
                        page_text, avg_conf, low_conf, line_conf = ocr_image(
                            image, lang=lang, advanced=False, fast=False
                        )
                        warnings.append(f"Fallback OCR used for {file.filename}.")
                    except Exception as exc2:
                        if not fast_ocr:
                            try:
                                page_text, avg_conf, low_conf, line_conf = ocr_image(
                                    image, lang=lang, advanced=True, fast=False
                                )
                                warnings.append(f"Fallback OCR used for {file.filename}.")
                            except Exception as exc3:
                                warnings.append(f"OCR failed on {file.filename}: {exc3}")
                                continue
                        else:
                            warnings.append(f"OCR failed on {file.filename}: {exc2}")
                            continue
                else:
                    warnings.append(f"OCR failed on {file.filename}: {exc}")
                    continue
            if fast_ocr and not advanced_ocr:
                short_text = len((page_text or "").strip()) < 12
                low_conf = avg_conf and avg_conf < 40
                if short_text or low_conf:
                    try:
                        page_text, avg_conf, low_conf, line_conf = ocr_image(
                            image, lang=lang, advanced=False, fast=False
                        )
                        if (page_text or "").strip():
                            warnings.append(f"Fallback OCR used for {file.filename}.")
                    except Exception:
                        pass

            if not (page_text or "").strip() and fast_ocr and not advanced_ocr:
                try:
                    page_text, avg_conf, low_conf, line_conf = ocr_image(
                        image, lang=lang, advanced=False, fast=False
                    )
                    if (page_text or "").strip():
                        warnings.append(f"Fallback OCR used for {file.filename}.")
                except Exception:
                    pass
            if not (page_text or "").strip() and not advanced_ocr and not fast_ocr:
                try:
                    page_text, avg_conf, low_conf, line_conf = ocr_image(
                        image, lang=lang, advanced=True, fast=False
                    )
                    if (page_text or "").strip():
                        warnings.append(f"Fallback OCR used for {file.filename}.")
                except Exception:
                    pass

            _add_page(page_text, avg_conf=avg_conf, low_conf=low_conf, line_conf=line_conf)

            if ext in image_exts:
                if skip_storage:
                    if privacy_mode and _supabase_only():
                        try:
                            os.remove(file_path)
                        except Exception:
                            pass
                        return None, warnings, "Privacy mode requires local storage, but local is disabled."
                    if privacy_mode or not store_local:
                        try:
                            os.remove(file_path)
                        except Exception:
                            pass
                    continue
                storage_path = local_rel_path.replace("\\", "/")
                upload_ok = False
                if storage_available:
                    try:
                        upload_ok = upload_to_storage(file_path, storage_path) is not None
                    except Exception as exc:
                        warnings.append(f"Storage upload failed; kept locally: {exc}")
                        upload_ok = False

                if not privacy_mode:
                    image_paths.append(storage_path)
                elif _supabase_only():
                    try:
                        os.remove(file_path)
                    except Exception:
                        pass
                    return None, warnings, "Privacy mode requires local storage, but local is disabled."

                # Remove local copy when privacy mode, or when safely uploaded to cloud
                delete_local = False
                if privacy_mode:
                    delete_local = True
                elif storage_available and upload_ok and not store_local:
                    delete_local = True

                if delete_local:
                    try:
                        os.remove(file_path)
                    except Exception:
                        pass

        if not combined_text:
            if warnings:
                return None, warnings, f"No valid images were processed. {warnings[0]}"
            if empty_pages:
                return (
                    None,
                    warnings,
                    "OCR ran but no readable text was found. Try a clearer image or check OCR setup.",
                )
            return (
                None,
                warnings,
                "No valid files were processed. Please upload images (JPG/PNG/HEIC), PDF, DOCX, or TXT.",
            )

        extracted_text = "\n\n".join(combined_text).strip()
        cleaned_text = clean_text(extracted_text) if cleanup else extracted_text
        intent = detect_intent(cleaned_text if detect_intent_flag else extracted_text) if detect_intent_flag else "auto"
        mode = "auto"
        if intent == "document":
            mode = "pdf_doc"
        elif intent == "notes":
            mode = "study"
        elif intent == "form":
            mode = "autofill"
        elif intent == "quote":
            mode = "creative"
        language = detect_language(cleaned_text)
        if autocorrect:
            try:
                cleaned_text = auto_correct_text(cleaned_text, language)
            except Exception as exc:
                warnings.append(f"Auto-correct skipped: {exc}")

        summary = ""
        key_points = []
        mcqs = []
        if student_mode:
            try:
                summary, key_points, mcqs = student_pack(cleaned_text)
            except Exception as exc:
                warnings.append(f"Student mode skipped: {exc}")

        confidence_avg = round(sum(combined_conf) / len(combined_conf), 2) if combined_conf else 0.0

        scan = {
            "id": scan_id,
            "user_id": user_id,
            "image_paths": image_paths,
            "extracted_text": extracted_text,
            "cleaned_text": cleaned_text,
            "language": language,
            "intent": intent,
            "mode": mode,
            "confidence_avg": confidence_avg,
            "low_confidence_words": combined_low_conf[:50],
            "line_confidence": combined_line_conf,
            "summary": summary,
            "key_points": key_points,
            "mcqs": mcqs,
            "tags": [],
            "is_private": privacy_mode,
            "created_at": now_iso(),
        }

        try:
            upsert_scan(scan)
        except Exception as exc:
            warnings.append(f"Database save failed: {exc}")
        if use_flash:
            for warn in warnings:
                flash(warn, "warn")
        return scan, warnings, None

    @app.route("/dashboard")
    def dashboard():
        if not require_login():
            return redirect(url_for("login"))

        user_id = get_user_id()
        scan_id = request.args.get("scan_id")
        latest_scan = None
        try:
            scans = list_scans(user_id)
        except Exception as exc:
            flash(f"Database error: {exc}", "warn")
            scans = []
        if scan_id:
            try:
                latest_scan = get_scan(scan_id)
                if latest_scan and latest_scan.get("user_id") != user_id:
                    latest_scan = None
            except Exception:
                latest_scan = None
        recent_scans = scans[:5]
        tags = sorted({tag for scan in scans for tag in scan.get("tags", [])})

        return render_template(
            "dashboard.html",
            recent_scans=recent_scans,
            tags=tags,
            user_email=session.get("user_email"),
            latest_scan=latest_scan,
        )

    @app.route("/upload", methods=["POST"])
    def upload():
        if not require_login():
            return redirect(url_for("login"))

        user_id = get_user_id()
        try:
            scan, warnings, error = _process_upload(
                request.files.getlist("images"), user_id, use_flash=True
            )
        except Exception as exc:
            current_app.logger.exception("OCR upload failed")
            flash(f"OCR failed: {exc}", "warn")
            return redirect(url_for("dashboard"))
        if error:
            flash(error, "warn")
            return redirect(url_for("dashboard"))

        flash("OCR complete.", "success")
        return redirect(url_for("dashboard", scan_id=scan["id"]))

    @app.route("/api/health")
    def api_health():
        return jsonify({"ok": True})

    @app.route("/api/scans", methods=["GET"])
    def api_scans():
        user_id = _api_user_id()
        if not user_id:
            return jsonify({"ok": False, "error": "Unauthorized"}), 401
        try:
            scans = list_scans(user_id)
            return jsonify({"ok": True, "data": scans})
        except Exception as exc:
            return jsonify({"ok": False, "error": str(exc)}), 500

    @app.route("/api/scans", methods=["POST"])
    def api_upload():
        user_id = _api_user_id()
        if not user_id:
            return jsonify({"ok": False, "error": "Unauthorized"}), 401
        try:
            scan, warnings, error = _process_upload(
                request.files.getlist("images"), user_id, use_flash=False
            )
        except Exception as exc:
            current_app.logger.exception("API OCR failed")
            return jsonify({"ok": False, "error": f"OCR failed: {exc}"}), 500
        if error:
            return jsonify({"ok": False, "error": error, "warnings": warnings}), 400
        return jsonify({"ok": True, "data": scan, "warnings": warnings})

    @app.route("/api/scans/<scan_id>", methods=["GET"])
    def api_scan(scan_id):
        user_id = _api_user_id()
        if not user_id:
            return jsonify({"ok": False, "error": "Unauthorized"}), 401
        try:
            scan = get_scan(scan_id)
        except Exception as exc:
            return jsonify({"ok": False, "error": str(exc)}), 500
        if not scan:
            return jsonify({"ok": False, "error": "Scan not found"}), 404
        if scan.get("user_id") != user_id:
            return jsonify({"ok": False, "error": "Forbidden"}), 403
        return jsonify({"ok": True, "data": scan})

    @app.route("/api/scans/<scan_id>/translate", methods=["POST"])
    def api_scan_translate(scan_id):
        user_id = _api_user_id()
        if not user_id:
            return jsonify({"ok": False, "error": "Unauthorized"}), 401
        try:
            scan = get_scan(scan_id)
        except Exception as exc:
            return jsonify({"ok": False, "error": str(exc)}), 500
        if not scan:
            return jsonify({"ok": False, "error": "Scan not found"}), 404
        if scan.get("user_id") != user_id:
            return jsonify({"ok": False, "error": "Forbidden"}), 403

        payload = request.get_json(silent=True) or {}
        target_lang = (payload.get("target_lang") or "en").lower()
        source_text = scan.get("cleaned_text") or scan.get("extracted_text", "")
        if not source_text:
            return jsonify({"ok": False, "error": "No text available"}), 400

        try:
            from services.free_ai import translate_text

            translated, source_lang = translate_text(source_text, target_lang)
            scan["translation"] = {
                "target": target_lang,
                "source": source_lang,
                "text": translated,
                "created_at": now_iso(),
            }
            upsert_scan(scan)
            log_translation(scan_id, scan.get("user_id"), source_lang, target_lang, translated)
            return jsonify(
                {
                    "ok": True,
                    "data": {
                        "translation_text": translated,
                        "target": target_lang,
                        "source": source_lang,
                    },
                }
            )
        except Exception as exc:
            return jsonify({"ok": False, "error": f"Translation failed: {exc}"}), 500

    @app.route("/api/scans/<scan_id>", methods=["PATCH"])
    def api_scan_update(scan_id):
        user_id = _api_user_id()
        if not user_id:
            return jsonify({"ok": False, "error": "Unauthorized"}), 401
        try:
            scan = get_scan(scan_id)
        except Exception as exc:
            return jsonify({"ok": False, "error": str(exc)}), 500
        if not scan:
            return jsonify({"ok": False, "error": "Scan not found"}), 404
        if scan.get("user_id") != user_id:
            return jsonify({"ok": False, "error": "Forbidden"}), 403

        payload = request.get_json(silent=True) or {}
        cleaned_text = (payload.get("cleaned_text") or scan.get("cleaned_text") or "").strip()
        tags = payload.get("tags") or scan.get("tags") or []
        if isinstance(tags, str):
            tags = [tag.strip() for tag in tags.split(",") if tag.strip()]

        scan["cleaned_text"] = cleaned_text
        scan["tags"] = tags
        scan["updated_at"] = now_iso()
        try:
            upsert_scan(scan)
        except Exception as exc:
            return jsonify({"ok": False, "error": str(exc)}), 500
        return jsonify({"ok": True, "data": scan})

    @app.route("/api/scans/<scan_id>", methods=["DELETE"])
    def api_scan_delete(scan_id):
        user_id = _api_user_id()
        if not user_id:
            return jsonify({"ok": False, "error": "Unauthorized"}), 401
        try:
            scan = get_scan(scan_id)
        except Exception as exc:
            return jsonify({"ok": False, "error": str(exc)}), 500
        if not scan:
            return jsonify({"ok": False, "error": "Scan not found"}), 404
        if scan.get("user_id") != user_id:
            return jsonify({"ok": False, "error": "Forbidden"}), 403

        image_paths = scan.get("image_paths", [])
        for path in image_paths:
            local_path = os.path.join(current_app.config["UPLOAD_FOLDER"], path)
            if os.path.exists(local_path):
                try:
                    os.remove(local_path)
                except Exception:
                    pass
        try:
            delete_from_storage(image_paths)
        except Exception:
            pass
        try:
            delete_scan_record(scan_id)
        except Exception as exc:
            return jsonify({"ok": False, "error": str(exc)}), 500
        return jsonify({"ok": True})

    @app.route("/result/<scan_id>")
    def result(scan_id):
        if not require_login():
            return redirect(url_for("login"))

        try:
            scan = get_scan(scan_id)
        except Exception as exc:
            flash(f"Database error: {exc}", "warn")
            return redirect(url_for("dashboard"))
        if not scan:
            flash("Scan not found.", "warn")
            return redirect(url_for("dashboard"))

        actions = extract_actions(scan.get("cleaned_text") or scan.get("extracted_text", ""))

        extracted = scan.get("extracted_text") or ""
        word_count = len([w for w in extracted.split() if w.strip()])
        file_name = ""
        if scan.get("image_paths"):
            file_name = os.path.basename(scan["image_paths"][0])
        source_label = "Upload"
        if file_name.lower().endswith("camera.jpg"):
            source_label = "Camera"

        scan["word_count"] = word_count
        scan["file_name"] = file_name or "N/A"
        scan["source_label"] = source_label
        scan["extracted_excerpt"] = extracted[:500]

        return render_template(
            "result.html",
            scan=scan,
            actions=actions,
        )

    @app.route("/result/<scan_id>/save", methods=["POST"])
    def save_result(scan_id):
        if not require_login():
            return redirect(url_for("login"))

        try:
            scan = get_scan(scan_id)
        except Exception as exc:
            flash(f"Database error: {exc}", "warn")
            return redirect(url_for("dashboard"))
        if not scan:
            flash("Scan not found.", "warn")
            return redirect(url_for("dashboard"))

        cleaned_text = request.form.get("cleaned_text", "").strip()
        tags = [tag.strip() for tag in request.form.get("tags", "").split(",") if tag.strip()]

        scan["cleaned_text"] = cleaned_text
        scan["tags"] = tags
        scan["updated_at"] = now_iso()

        try:
            upsert_scan(scan)
        except Exception as exc:
            flash(f"Save failed: {exc}", "warn")
            return redirect(url_for("result", scan_id=scan_id))
        flash("Saved successfully.", "success")
        return redirect(url_for("result", scan_id=scan_id))

    @app.route("/result/<scan_id>/export", methods=["POST"])
    def export_result(scan_id):
        if not require_login():
            return redirect(url_for("login"))

        try:
            scan = get_scan(scan_id)
        except Exception as exc:
            flash(f"Database error: {exc}", "warn")
            return redirect(url_for("dashboard"))
        if not scan:
            flash("Scan not found.", "warn")
            return redirect(url_for("dashboard"))

        export_format = request.form.get("export_format", "txt")
        content = scan.get("cleaned_text") or scan.get("extracted_text", "")

        if export_format == "txt":
            buffer = io.BytesIO(content.encode("utf-8"))
            filename = f"visiontext_{scan_id}.txt"
            log_export(scan_id, scan.get("user_id"), "txt")
            return send_file(buffer, as_attachment=True, download_name=filename, mimetype="text/plain")

        if export_format == "pdf":
            try:
                from fpdf import FPDF

                pdf = FPDF()
                pdf.add_page()
                pdf.set_font("Arial", size=12)
                for line in content.split("\n"):
                    pdf.multi_cell(0, 8, line)
                pdf_bytes = pdf.output(dest="S").encode("latin-1")
                buffer = io.BytesIO(pdf_bytes)
                filename = f"visiontext_{scan_id}.pdf"
                log_export(scan_id, scan.get("user_id"), "pdf")
                return send_file(
                    buffer, as_attachment=True, download_name=filename, mimetype="application/pdf"
                )
            except Exception:
                flash("PDF export needs fpdf2 installed.", "warn")
                return redirect(url_for("result", scan_id=scan_id))

        if export_format == "docx":
            try:
                from docx import Document

                doc = Document()
                for line in content.split("\n"):
                    doc.add_paragraph(line)
                buffer = io.BytesIO()
                doc.save(buffer)
                buffer.seek(0)
                filename = f"visiontext_{scan_id}.docx"
                log_export(scan_id, scan.get("user_id"), "docx")
                return send_file(
                    buffer,
                    as_attachment=True,
                    download_name=filename,
                    mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                )
            except Exception:
                flash("DOCX export needs python-docx installed.", "warn")
                return redirect(url_for("result", scan_id=scan_id))

        flash("Unsupported export format.", "warn")
        return redirect(url_for("result", scan_id=scan_id))

    @app.route("/result/<scan_id>/translate", methods=["POST"])
    def translate_result(scan_id):
        if not require_login():
            return redirect(url_for("login"))

        try:
            scan = get_scan(scan_id)
        except Exception as exc:
            flash(f"Database error: {exc}", "warn")
            return redirect(url_for("dashboard"))
        if not scan:
            flash("Scan not found.", "warn")
            return redirect(url_for("dashboard"))

        target_lang = request.form.get("target_lang", "en")
        source_text = scan.get("cleaned_text") or scan.get("extracted_text", "")
        try:
            from services.free_ai import translate_text

            translated, source_lang = translate_text(source_text, target_lang)
            scan["translation"] = {
                "target": target_lang,
                "source": source_lang,
                "text": translated,
                "created_at": now_iso(),
            }
            upsert_scan(scan)
            log_translation(scan_id, scan.get("user_id"), source_lang, target_lang, translated)
            flash("Translation complete.", "success")
        except Exception as exc:
            flash(f"Translation failed: {exc}", "warn")
        return redirect(url_for("result", scan_id=scan_id))

    @app.route("/result/<scan_id>/tts", methods=["POST"])
    def tts_result(scan_id):
        if not require_login():
            return redirect(url_for("login"))
        try:
            scan = get_scan(scan_id)
        except Exception as exc:
            flash(f"Database error: {exc}", "warn")
            return redirect(url_for("dashboard"))
        if not scan:
            flash("Scan not found.", "warn")
            return redirect(url_for("dashboard"))

        text = scan.get("cleaned_text") or scan.get("extracted_text", "")
        if not text:
            flash("No text available for audio.", "warn")
            return redirect(url_for("result", scan_id=scan_id))

        lang_map = {
            "en": "en",
            "hi": "hi",
            "gu": "gu",
        }
        language_code = lang_map.get((scan.get("language") or "en")[:2], "en")
        speed = request.form.get("tts_speed", "normal")
        slow = speed == "slow"
        try:
            from services.free_ai import synthesize_speech

            audio_bytes = synthesize_speech(text[:4000], language_code=language_code, slow=slow)
            audio_name = f"audio_{scan_id}.mp3"
            storage_path = f"audio/{scan.get('user_id')}/{audio_name}"
            temp_path = os.path.join(current_app.config["UPLOAD_FOLDER"], storage_path)
            os.makedirs(os.path.dirname(temp_path), exist_ok=True)
            with open(temp_path, "wb") as handle:
                handle.write(audio_bytes)

            storage_env_on = os.getenv("SUPABASE_USE_STORAGE", "true").lower() not in {"false", "0", "no"}
            storage_available = False
            if storage_env_on:
                try:
                    require_storage_client()
                    storage_available = supabase_storage_client() is not None
                except Exception:
                    storage_available = False

            if storage_available:
                try:
                    upload_to_storage(temp_path, storage_path)
                    # safe to delete local copy
                    try:
                        os.remove(temp_path)
                    except Exception:
                        pass
                except Exception as exc:
                    flash(f"Audio upload failed; kept locally: {exc}", "warn")
            elif _supabase_only():
                try:
                    os.remove(temp_path)
                except Exception:
                    pass
                flash("Storage not available; audio not saved.", "warn")
                return redirect(url_for("result", scan_id=scan_id))

            # Always keep a path that /uploads can serve
            scan["audio_path"] = storage_path
            scan["audio_updated_at"] = now_iso()
            upsert_scan(scan)
            flash("Audio generated.", "success")
        except Exception as exc:
            flash(f"TTS failed: {exc}", "warn")
        return redirect(url_for("result", scan_id=scan_id))

    @app.route("/uploads/<path:filename>")
    def uploaded_file(filename):
        local_path = os.path.join(current_app.config["UPLOAD_FOLDER"], filename)
        if os.path.exists(local_path):
            return send_from_directory(current_app.config["UPLOAD_FOLDER"], filename)

        try:
            data, content_type = download_from_storage(filename)
            if data is not None:
                return send_file(io.BytesIO(data), mimetype=content_type)
        except Exception:
            pass

        return abort(404)
